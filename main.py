from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.responses import JSONResponse
from pydantic import BaseModel
import requests, io, re
from collections import Counter

# --- مستخرِجات النص ---
from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation
from docx import Document

app = FastAPI(title="StudyCoach Backend", version="2.0")

# ========================
# نماذج البيانات (لـ iOS)
# ========================
class SummaryBlock(BaseModel):
    title: str
    bullets: list[str]

class MCQ(BaseModel):
    q: str
    choices: list[str]
    answer: int
    explain: str | None = None

class AIResponse(BaseModel):
    ok: bool
    summary_blocks: list[SummaryBlock]
    mcq: list[MCQ]

class PDFRequest(BaseModel):
    url: str  # قد يكون PDF أو PPTX أو DOCX

# ========================
# أدوات مساعدة
# ========================
def clean_text(t: str) -> str:
    return re.sub(r"\s+", " ", (t or "")).strip()

def extract_text_from_pdf_bytes(b: bytes) -> str:
    with io.BytesIO(b) as bio:
        return pdf_extract_text(bio) or ""

def extract_text_from_pptx_bytes(b: bytes) -> str:
    prs = Presentation(io.BytesIO(b))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if txt:
                    chunks.append(txt)
    return "\n".join(chunks)

def extract_text_from_docx_bytes(b: bytes) -> str:
    doc = Document(io.BytesIO(b))
    chunks = []
    for p in doc.paragraphs:
        txt = (p.text or "").strip()
        if txt:
            chunks.append(txt)
    return "\n".join(chunks)

def naive_sentences(text: str) -> list[str]:
    # تقسيم بسيط مع دعم ؟ العربية
    parts = re.split(r'(?<=[\.\!\?\u061F])\s+', text)
    return [s.strip() for s in parts if len(s.strip()) > 30]

def top_terms(text: str, k=40) -> list[str]:
    words = re.findall(r"[A-Za-z\u0600-\u06FF][A-Za-z0-9_\u0600-\u06FF'-]*", text)
    words = [w.lower() for w in words if len(w) > 3]
    stop = {
        "this","that","with","from","your","about","these","those","which","their","have","will","there","here","where",
        "when","then","into","over","under","also","such","than","only","very","more","most","much","many","some",
        "been","were","what","shall","should","could","would","might","must","also","upon","into","between","within"
    }
    common = [w for w in words if w not in stop]
    return [w for w,_ in Counter(common).most_common(k)]

# ========================
# توليد الملخّص + MCQ (هيوريستك محلي)
# ========================
def build_summary_blocks(text: str, max_blocks: int = 4, bullets_per_block: int = 4) -> list[SummaryBlock]:
    sents = naive_sentences(text)
    if not sents:
        return [SummaryBlock(title="ملحوظة", bullets=["تعذّر استخراج جمل كافية من الملف."])]
    terms = set(top_terms(text, k=30))
    scored = []
    for s in sents:
        score = sum(1 for w in re.findall(r"[A-Za-z\u0600-\u06FF][A-Za-z0-9_\u0600-\u06FF'-]*", s.lower()) if w in terms)
        scored.append((score, s))
    # اختَر أفضل الجمل ورتّبها على شكل مجموعات
    top = [s for _, s in sorted(scored, key=lambda x: x[0], reverse=True)[:max_blocks * bullets_per_block]]
    blocks = []
    for i in range(0, len(top), bullets_per_block):
        chunk = top[i:i+bullets_per_block]
        if not chunk: break
        title = f"أفكار رئيسية #{len(blocks)+1}"
        bullets = [clean_text(s) for s in chunk]
        blocks.append(SummaryBlock(title=title, bullets=bullets))
        if len(blocks) >= max_blocks: break
    return blocks

def build_mcq(text: str, n: int = 6) -> list[MCQ]:
    """
    MCQ بسيط: نأخذ مصطلحات متكررة كإجابات صحيحة، ونولّد مشتتات من مصطلحات أخرى.
    هذا حل محلي سريع؛ لاحقًا ممكن نستبدله باستدعاء OpenAI.
    """
    terms = top_terms(text, k=50)
    if len(terms) < 6:
        # fallback بسيط
        return [
            MCQ(q="سؤال تجريبي: ما هي أفضل طريقة لاستخلاص الأفكار الرئيسية؟",
                choices=["القراءة السريعة", "التلخيص الآلي", "مراجعة الأسئلة", "كل ما سبق"],
                answer=3,
                explain="الدمج بين أكثر من طريقة يرفع جودة الاستيعاب.")
        ]

    sents = naive_sentences(text)
    qs = []
    used = set()
    i = 0
    for t in terms:
        if t in used: 
            continue
        # ابحث جملة تحتوي t
        sent = next((s for s in sents if re.search(rf"\b{re.escape(t)}\b", s, re.IGNORECASE)), None)
        if not sent:
            continue
        # خيارات: الصحيحة t + 3 مشتتات
        distractors = [d for d in terms if d != t][:12]  # مجموعة أولية
        import random
        random.shuffle(distractors)
        distractors = distractors[:3]
        choices = [t] + distractors
        random.shuffle(choices)
        answer_idx = choices.index(t)
        qs.append(MCQ(
            q=f"اختر المصطلح الذي يكمّل الفكرة: «{clean_text(sent)}»",
            choices=choices,
            answer=answer_idx,
            explain=f"المصطلح المناسب في السياق هو «{t}»."
        ))
        used.add(t)
        i += 1
        if i >= n:
            break
    if not qs:
        # fallback لو ما لقينا مطابقات
        qs = [MCQ(q="سؤال عام: أي مما يلي يُعد فائدة للتلخيص؟",
                  choices=["تقليل الوقت", "فهم أعمق", "تثبيت المعلومة", "كل ما سبق"],
                  answer=3,
                  explain="التلخيص الجيد يجمع كل ما سبق.")]
    return qs

def make_response_from_text(text: str) -> AIResponse:
    text = clean_text(text)
    if len(text) < 80:
        raise HTTPException(status_code=422, detail="تعذّر استخراج نص كافٍ من الملف.")
    blocks = build_summary_blocks(text)
    mcq = build_mcq(text)
    return AIResponse(ok=True, summary_blocks=blocks, mcq=mcq)

# ========================
# صحّة الخادم
# ========================
@app.get("/health")
def health():
    return {"ok": True}

# ========================
# من رابط (PDF/PPTX/DOCX)
# ========================
@app.post("/process_pdf", response_model=AIResponse)
def process_pdf(req: PDFRequest):
    try:
        r = requests.get(req.url, timeout=40)
        r.raise_for_status()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Download error: {e}")

    content_type = (r.headers.get("Content-Type") or "").lower()
    url_lower = req.url.lower()

    try:
        if "pdf" in content_type or url_lower.endswith(".pdf"):
            text = extract_text_from_pdf_bytes(r.content)
        elif "presentation" in content_type or url_lower.endswith(".pptx"):
            text = extract_text_from_pptx_bytes(r.content)
        elif "word" in content_type or url_lower.endswith(".docx"):
            text = extract_text_from_docx_bytes(r.content)
        else:
            # جرّب PDF كـ fallback
            try:
                text = extract_text_from_pdf_bytes(r.content)
            except Exception:
                raise HTTPException(status_code=415, detail=f"Unsupported file type: {content_type or 'unknown'}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Extract error: {e}")

    return make_response_from_text(text)

# ========================
# رفع ملف (PDF/PPTX/DOCX)
# ========================
@app.post("/process_pdf_upload", response_model=AIResponse)
async def process_pdf_upload(file: UploadFile = File(...)):
    try:
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="ملف فارغ")

        name = (file.filename or "").lower()
        if name.endswith(".pdf"):
            text = extract_text_from_pdf_bytes(content)
        elif name.endswith(".pptx"):
            text = extract_text_from_pptx_bytes(content)
        elif name.endswith(".docx"):
            text = extract_text_from_docx_bytes(content)
        else:
            raise HTTPException(status_code=415, detail="النوع غير مدعوم. ارفع PDF أو PPTX أو DOCX")

        return make_response_from_text(text)

    except HTTPException:
        raise
    except Exception as e:
        return JSONResponse(status_code=500, content={"ok": False, "error": f"Server error: {e.__class__.__name__}"})

# تشغيل محلي فقط
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
# =============== Moodle Proxy ===============
from fastapi import FastAPI, HTTPException, UploadFile, File, Query
from pydantic import BaseModel
import requests, io

class MoodCoursesReq(BaseModel):
    base: str   # مثل https://lms.yu.edu.sa
    token: str  # moodle_mobile_app token

@app.post("/moodle/courses")
def moodle_courses(req: MoodCoursesReq):
    # core_enrol_get_users_courses (يتطلب user id عادةً؛ لكن مع توكن الموبايل نقدر نجيب أنا ذاتياً)
    # نستخدم خدمة tool_mobile_call_external_functions عبر server.php
    url = f"{req.base}/webservice/rest/server.php"
    payload = {
        "wstoken": req.token,
        "moodlewsrestformat": "json",
        "wsfunction": "core_enrol_get_users_courses",
        "userid": 0  # 0 = المستخدم الحالي للتوكن في بعض تركيبات الموودل، وإن فشل نرجع رسالة
    }
    r = requests.post(url, data=payload, timeout=40)
    if r.status_code != 200:
        raise HTTPException(502, "Moodle courses fetch failed")
    data = r.json()
    if isinstance(data, dict) and data.get("exception"):
        # محاولة بديلة: tool_mobile_get_course_contents مع استدعاء آخر يتطلب courseid لاحقاً
        raise HTTPException(400, f"Moodle error: {data.get('message','unknown')}")
    # نبسّط: id, shortname, fullname
    out = [{"id": c.get("id"), "shortname": c.get("shortname"), "fullname": c.get("fullname")} for c in data]
    return {"ok": True, "courses": out}

class MoodContentsReq(BaseModel):
    base: str
    token: str
    courseid: int

@app.post("/moodle/contents")
def moodle_course_contents(req: MoodContentsReq):
    # core_course_get_contents لإحضار ملفات المقرر
    url = f"{req.base}/webservice/rest/server.php"
    payload = {
        "wstoken": req.token,
        "moodlewsrestformat": "json",
        "wsfunction": "core_course_get_contents",
        "courseid": req.courseid
    }
    r = requests.post(url, data=payload, timeout=40)
    if r.status_code != 200:
        raise HTTPException(502, "Moodle contents fetch failed")
    data = r.json()
    if isinstance(data, dict) and data.get("exception"):
        raise HTTPException(400, f"Moodle error: {data.get('message','unknown')}")
    # نستخرج الملفات القابلة للتنزيل (resource, fileurl)
    files = []
    for section in data:
        for mod in section.get("modules", []):
            for f in mod.get("contents", []) or []:
                fileurl = f.get("fileurl")
                if fileurl:
                    files.append({
                        "name": f.get("filename") or f.get("filepath") or mod.get("name"),
                        "fileurl": fileurl,  # سنستخدمه عبر بروكسي التحميل أدناه
                        "modname": mod.get("modname"),
                        "section": section.get("name") or ""
                    })
    return {"ok": True, "files": files}

@app.get("/moodle/fetch")
def moodle_fetch(fileurl: str = Query(...), token: str = Query(...)):
    # تنزيل الملف من Moodle باستخدام توكن (fileurl+token=...)
    if "token=" not in fileurl:
        sep = "&" if "?" in fileurl else "?"
        fileurl = f"{fileurl}{sep}token={token}"
    rr = requests.get(fileurl, timeout=60)
    if rr.status_code != 200:
        raise HTTPException(400, "Failed to fetch file from Moodle")
    # نعيد البايتات كما هي (سيستقبلها الـiOS عبر endpoint آخر أو نستعملها داخليًا)
    return rr.content
